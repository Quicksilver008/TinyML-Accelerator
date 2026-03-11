"""
EdgeMATX Presentation Generator  v2
Produces: EdgeMATX_presentation.pptx  (17 slides)
Run:  python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from lxml import etree
import os

# ── THEME ─────────────────────────────────────────────────────────────────────
BG     = RGBColor(0x0B, 0x14, 0x2A)
CARD   = RGBColor(0x14, 0x21, 0x3E)
CARD2  = RGBColor(0x1E, 0x2D, 0x4F)
CYAN   = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GREEN  = RGBColor(0x00, 0xDC, 0x75)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x90, 0xA4, 0xB8)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
PINK   = RGBColor(0xF0, 0x6A, 0xC8)
RED    = RGBColor(0xFF, 0x4D, 0x4D)

W = Inches(13.333)
H = Inches(7.5)

# ── CORE HELPERS ──────────────────────────────────────────────────────────────

def make_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, text, l, t, w, h, size=22, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_para(tf, text, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if before:
        p.space_before = Pt(before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p

def rect(slide, l, t, w, h, fill=CARD, line=None, lw=Pt(1.5), radius=False):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape

def hrule(slide, l, t, w, color=CYAN, thick=Pt(2)):
    s = slide.shapes.add_shape(1, l, t, w, thick)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def header(slide, title, sub=None):
    box(slide, title,
        Inches(0.45), Inches(0.12), Inches(12.5), Inches(0.75),
        size=34, bold=True, color=CYAN)
    hrule(slide, Inches(0.45), Inches(0.9), Inches(12.4))
    if sub:
        box(slide, sub,
            Inches(0.45), Inches(0.92), Inches(12.4), Inches(0.45),
            size=15, color=GRAY)

def footer(slide, text="EdgeMATX  |  NITK Surathkal  |  ECE Department  |  2026"):
    box(slide, text,
        Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.32),
        size=11, color=GRAY)

def pill(slide, l, t, label, fill=CYAN, text_col=BG):
    """Small coloured badge/pill."""
    r = rect(slide, l, t, Inches(1.8), Inches(0.42), fill=fill)
    box(slide, label, l, t + Inches(0.07),
        Inches(1.8), Inches(0.35), size=13, bold=True,
        color=text_col, align=PP_ALIGN.CENTER)
    return r

def img(slide, fname, l, t, w, h, subdir="figures"):
    """Insert an image if file exists; show a labelled placeholder otherwise."""
    base = os.path.dirname(os.path.abspath(__file__))
    path = os.path.join(base, subdir, fname) if subdir else os.path.join(base, fname)
    if os.path.exists(path):
        return slide.shapes.add_picture(path, l, t, w, h)
    # Fallback placeholder
    r = rect(slide, l, t, w, h, fill=RGBColor(0x08, 0x10, 0x1C), line=GRAY)
    box(slide, f"[ {fname} ]", l, t + Inches(0.1), w, Inches(0.4),
        size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    return r

def url_box(slide, text, url, l, t, w, h, size=13, col=CYAN):
    """Text box whose first run carries a clickable hyperlink."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = col
    run.font.underline = True
    try:
        run.hyperlink.address = url
    except Exception:
        pass
    return tb

# ── ANIMATION HELPER ──────────────────────────────────────────────────────────

def add_click_animations(slide, shape_ids, preset=10, delay_ms=0):
    """Add on-click fade-in entrance animations to a list of shape IDs."""
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    base_id = 3

    seq_children = ""
    for i, spid in enumerate(shape_ids):
        oid = base_id + i * 5
        seq_children += f"""
        <p:par xmlns:p="{p_ns}">
          <p:cTn id="{oid}" fill="hold">
            <p:stCondLst><p:cond evt="onBegin" delay="indefinite"/></p:stCondLst>
            <p:childTnLst>
              <p:par>
                <p:cTn id="{oid+1}" presetID="{preset}" presetClass="entr"
                       presetSubtype="0" fill="hold" grpId="{i}"
                       nodeType="clickEffect">
                  <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
                  <p:childTnLst>
                    <p:set>
                      <p:cBhvr>
                        <p:cTn id="{oid+2}" dur="500" fill="hold"/>
                        <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                        <p:attrNameLst>
                          <p:attrName>style.visibility</p:attrName>
                        </p:attrNameLst>
                      </p:cBhvr>
                      <p:to><p:strVal val="visible"/></p:to>
                    </p:set>
                  </p:childTnLst>
                </p:cTn>
              </p:par>
            </p:childTnLst>
          </p:cTn>
        </p:par>"""

    xml = f"""<p:timing xmlns:p="{p_ns}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{seq_children}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onBegin" delay="0"><p:tn/></p:cond>
            </p:prevCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""
    slide.element.append(etree.fromstring(xml))

def add_transition(slide, dur_ms=700, t_type="fade"):
    """Add a slide transition."""
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    xml = (f'<p:transition xmlns:p="{p_ns}" spd="med" dur="{dur_ms}">'
           f'<p:{t_type}/></p:transition>')
    slide.element.append(etree.fromstring(xml))

# =============================================================================
# SLIDES
# =============================================================================

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank(prs)
    set_bg(s)

    # Left accent stripe
    rect(s, Inches(0), Inches(0), Inches(0.12), H, fill=CYAN)
    # Decorative top bar
    rect(s, Inches(0.12), Inches(0), W, Inches(0.06), fill=ORANGE)

    # Large project acronym
    box(s, "EdgeMATX",
        Inches(0.7), Inches(0.7), Inches(9), Inches(1.6),
        size=80, bold=True, color=CYAN)

    # Full title
    box(s, "Design and Integration of a TinyML Systolic\nAccelerator with PicoRV32",
        Inches(0.7), Inches(2.3), Inches(8.8), Inches(1.3),
        size=26, color=WHITE)

    hrule(s, Inches(0.7), Inches(3.65), Inches(8.5), color=ORANGE, thick=Pt(1.5))

    # Team block
    tb = s.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(6), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    add_para(tf, "Submitted by", 13, italic=True, color=GRAY)
    for name in ["Nishchay Pallav  221EC233",
                 "Mohammad Omar Sulemani  221EC230",
                 "Md Atib Kaif  221EC129"]:
        add_para(tf, name, 17, bold=True, color=WHITE, before=2)

    # Guide
    box(s, "Under the Guidance of  Dr Rathamala Rao",
        Inches(0.7), Inches(5.5), Inches(7), Inches(0.45),
        size=15, color=GRAY)

    # Institution
    box(s, "Dept. of Electronics & Communication Engineering\n"
           "National Institute of Technology Karnataka, Surathkal  575025",
        Inches(0.7), Inches(6.1), Inches(9), Inches(0.8),
        size=13, color=GRAY)

    # Right side: stats panel
    rect(s, Inches(10.0), Inches(1.5), Inches(3.0), Inches(4.5),
         fill=CARD, line=CYAN, lw=Pt(1))

    for i, (val, lbl, col) in enumerate([
        ("673",    "Cycles",         CYAN),
        ("54x",    "Peak Speedup",   GREEN),
        ("19/19",  "Tests Pass",     ORANGE),
        ("4x4",    "Systolic Array", YELLOW),
    ]):
        y = Inches(1.7) + i * Inches(1.05)
        box(s, val, Inches(10.15), y, Inches(2.7), Inches(0.58),
            size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
        box(s, lbl, Inches(10.15), y + Inches(0.52), Inches(2.7), Inches(0.38),
            size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # NITK emblem — top-right corner, clear of all panels
    img(s, "NITK_Emblem.png",
        Inches(12.0), Inches(0.08), Inches(1.15), Inches(1.15), subdir="")

    add_transition(s)
    return s

# ── SLIDE 2: INTRODUCTION ─────────────────────────────────────────────────────
def slide_introduction(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Introduction",
           "A hardware coprocessor for matrix multiplication on embedded RISC-V")
    footer(s)

    # Left: context text
    rect(s, Inches(0.4), Inches(1.35), Inches(5.8), Inches(5.7),
         fill=CARD, line=CYAN, lw=Pt(1.5))

    tb = s.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(5.55), Inches(5.4))
    tf = tb.text_frame; tf.word_wrap = True

    add_para(tf, "What is EdgeMATX?", 18, bold=True, color=CYAN)
    add_para(tf, "A custom hardware coprocessor tightly coupled to the open-source "
             "PicoRV32 RISC-V processor.  A single custom instruction replaces 64 "
             "sequential software operations with 16 fully parallel hardware MACs.",
             13, color=WHITE, before=5)

    add_para(tf, "Why Does It Matter?", 18, bold=True, color=ORANGE, before=12)
    for pt in [
        "TinyML inference on edge devices is dominated by dense matrix operations",
        "General-purpose RISC-V cores execute these sequentially -- very slowly",
        "Purpose-built systolic hardware delivers 38-54x speedup at minimal area",
    ]:
        add_para(tf, "  >  " + pt, 12, color=WHITE, before=4)

    add_para(tf, "Our Approach", 18, bold=True, color=GREEN, before=12)
    add_para(tf, "One opcode. CPU stalls. 4x4 systolic array fires. Results appear "
             "in memory. 673 deterministic cycles -- every time.",
             13, color=WHITE, before=5)

    add_para(tf, "Scope & Status", 18, bold=True, color=YELLOW, before=12)
    add_para(tf, "RTL simulation complete. 19/19 tests passing. "
             "Two live web visualizers. FPGA deployment upcoming.",
             13, color=WHITE, before=5)

    # Right: system stack image
    img(s, "fig_system_stack.png",
        Inches(6.45), Inches(1.35), Inches(6.5), Inches(5.7))

    add_transition(s)
    return s

# ── SLIDE 3: THE PROBLEM ──────────────────────────────────────────────────────
def slide_problem(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "The Problem",
           "Matrix multiplication is ubiquitous -- and painfully slow on embedded CPUs")
    footer(s)

    # Left column: where it's used
    rect(s, Inches(0.4), Inches(1.35), Inches(4.0), Inches(5.7),
         fill=CARD, line=CYAN, lw=Pt(1))
    box(s, "Where It's Needed",
        Inches(0.5), Inches(1.45), Inches(3.8), Inches(0.5),
        size=16, bold=True, color=CYAN)
    hrule(s, Inches(0.5), Inches(1.95), Inches(3.7), color=CYAN, thick=Pt(1))
    tb = s.shapes.add_textbox(Inches(0.55), Inches(2.05), Inches(3.75), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    for item in [
        ("  TinyML Inference",       "Neural network layers are dense GEMM ops"),
        ("  Digital Signal Proc.",   "Convolution = repeated matrix multiply"),
        ("  Scientific Computing",   "Linear solvers, transforms, all matrix ops"),
        ("  Image Processing",       "Filters, transforms, feature extraction"),
    ]:
        add_para(tf, item[0], 16, bold=True, color=WHITE, before=8)
        add_para(tf, item[1], 12, italic=True, color=GRAY, before=1)

    # Middle column: sequential execution problem
    rect(s, Inches(4.7), Inches(1.35), Inches(4.3), Inches(5.7),
         fill=CARD, line=ORANGE, lw=Pt(1))
    box(s, "Sequential CPU Execution",
        Inches(4.8), Inches(1.45), Inches(4.1), Inches(0.5),
        size=16, bold=True, color=ORANGE)
    hrule(s, Inches(4.8), Inches(1.95), Inches(4.1), color=ORANGE, thick=Pt(1))

    # Code snippet
    rect(s, Inches(4.8), Inches(2.05), Inches(4.1), Inches(2.3),
         fill=RGBColor(0x08, 0x0F, 0x1E), line=GRAY, lw=Pt(0.5))
    tb = s.shapes.add_textbox(Inches(4.9), Inches(2.1), Inches(3.9), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = False
    for line_txt in [
        "for i in range(4):",
        "  for j in range(4):",
        "    for k in range(4):",
        "      C[i][j] += A[i][k]",
        "             * B[k][j]",
    ]:
        add_para(tf, line_txt, 14, color=GREEN, align=PP_ALIGN.LEFT)

    tb2 = s.shapes.add_textbox(Inches(4.8), Inches(4.45), Inches(4.1), Inches(2.4))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for txt, col in [
        ("64 multiply-accumulate ops",    WHITE),
        ("executed one at a time",        GRAY),
        ("",                              WHITE),
        ("RV32I (no MUL):   26,130 cycles", RED),
        ("RV32IM (with MUL):  7,975 cycles", YELLOW),
    ]:
        add_para(tf2, txt, 15, color=col, before=3)

    # Right column: the gap
    rect(s, Inches(9.25), Inches(1.35), Inches(3.65), Inches(5.7),
         fill=CARD2, line=GREEN, lw=Pt(1))
    box(s, "The Gap",
        Inches(9.35), Inches(1.45), Inches(3.45), Inches(0.5),
        size=16, bold=True, color=GREEN)
    hrule(s, Inches(9.35), Inches(1.95), Inches(3.45), color=GREEN, thick=Pt(1))

    for i, (val, lbl, col) in enumerate([
        ("26 K",  "cycles wasted per\n4x4 multiply",   RED),
        ("100%",  "sequential\nexecution on CPU",       ORANGE),
        ("0x",    "hardware\nparallelism",              GRAY),
    ]):
        y = Inches(2.1) + i * Inches(1.35)
        box(s, val, Inches(9.35), y, Inches(3.45), Inches(0.65),
            size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
        box(s, lbl, Inches(9.35), y + Inches(0.6), Inches(3.45), Inches(0.6),
            size=13, color=GRAY, align=PP_ALIGN.CENTER)

    add_transition(s)
    return s

# ── SLIDE 4: THE SOLUTION ─────────────────────────────────────────────────────
def slide_solution(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "The Solution: EdgeMATX",
           "A custom PCPI coprocessor tightly coupled to PicoRV32")
    footer(s)

    # Big central tagline
    rect(s, Inches(0.4), Inches(1.35), Inches(12.5), Inches(1.2),
         fill=CARD, line=CYAN, lw=Pt(1.5))
    box(s, "One custom instruction replaces 64 sequential CPU operations",
        Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.0),
        size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Three pillars
    for i, (title, body, col) in enumerate([
        ("Tightly Coupled",
         "PCPI interface: zero memory-mapped overhead. CPU stalls and resumes "
         "-- no interrupts, no DMA.",
         CYAN),
        ("Systolic Datapath",
         "4x4 array of 16 processing elements. All 16 MACs execute in parallel "
         "every clock cycle.",
         ORANGE),
        ("Q5.10 Fixed-Point",
         "16-bit signed arithmetic. No floating-point hardware needed -- ideal "
         "for constrained silicon.",
         GREEN),
    ]):
        x = Inches(0.4) + i * Inches(4.32)
        rect(s, x, Inches(2.8), Inches(4.1), Inches(3.0),
             fill=CARD, line=col, lw=Pt(1.5))
        box(s, title,
            x + Inches(0.15), Inches(2.92), Inches(3.8), Inches(0.6),
            size=18, bold=True, color=col)
        hrule(s, x + Inches(0.15), Inches(3.52), Inches(3.8),
              color=col, thick=Pt(0.8))
        box(s, body,
            x + Inches(0.15), Inches(3.62), Inches(3.8), Inches(1.9),
            size=14, color=WHITE)

    # Bottom bar: key numbers
    rect(s, Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.85),
         fill=RGBColor(0x00, 0x1A, 0x33), line=CYAN, lw=Pt(1))
    for i, (val, lbl) in enumerate([
        ("673",     "Deterministic cycles"),
        ("38.83x",  "Speedup (sparse data)"),
        ("53.86x",  "Speedup (dense data)"),
        ("11.85x",  "Speedup vs MUL rv32im"),
    ]):
        x = Inches(0.6) + i * Inches(3.1)
        box(s, val, x, Inches(6.1), Inches(1.5), Inches(0.42),
            size=22, bold=True, color=CYAN)
        box(s, lbl, x + Inches(1.5), Inches(6.17), Inches(1.5), Inches(0.32),
            size=11, color=GRAY)

    add_transition(s)
    return s

# ── SLIDE 5: PROJECT PROGRESS ─────────────────────────────────────────────────
def slide_progress(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Project Progress",
           "Simulation-complete milestone -- all RTL objectives achieved")
    footer(s)

    phases = [
        ("DONE",  "Phase 1: RTL Design", [
            "4x4 systolic array in Verilog (Q5.10 fixed-point MACs)",
            "PCPI coprocessor wrapper for PicoRV32",
            "Memory-mapped I/O:  A@0x100  B@0x140  C@0x200",
        ], GREEN),
        ("DONE",  "Phase 2: Simulation & Verification", [
            "19/19 standalone accelerator tests -- 100% PASS",
            "8/8 CPU integration regression tests -- 100% PASS",
            "Cycle benchmarking: 673 deterministic cycles confirmed",
        ], CYAN),
        ("DONE",  "Phase 3: Tooling & Documentation", [
            "Two interactive web visualizers deployed on Vercel",
            "Full LaTeX project report submitted",
            "38-54x speedup measured and documented",
        ], ORANGE),
        ("NEXT",  "Phase 4: FPGA Deployment", [
            "Synthesize on Pynq-Z2 (Xilinx Zynq-7020)",
            "On-board cycle validation and timing closure",
            "AXI-BRAM memory interface integration",
        ], YELLOW),
    ]

    for i, (status, phase, items, col) in enumerate(phases):
        row = i // 2
        c   = i %  2
        x = Inches(0.4) + c * Inches(6.5)
        y = Inches(1.35) + row * Inches(2.95)

        rect(s, x, y, Inches(6.2), Inches(2.8),
             fill=CARD, line=col, lw=Pt(1.5))

        # Status badge
        badge_fill = GREEN if status == "DONE" else YELLOW
        rect(s, x + Inches(0.12), y + Inches(0.12),
             Inches(0.9), Inches(0.32), fill=badge_fill)
        box(s, status, x + Inches(0.12), y + Inches(0.12),
            Inches(0.9), Inches(0.32), size=10, bold=True,
            color=BG, align=PP_ALIGN.CENTER)

        box(s, phase, x + Inches(1.15), y + Inches(0.1),
            Inches(4.9), Inches(0.38), size=15, bold=True, color=col)
        hrule(s, x + Inches(0.12), y + Inches(0.54),
              Inches(5.95), color=col, thick=Pt(0.8))

        tb = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.68),
                                   Inches(5.9), Inches(1.9))
        tf = tb.text_frame; tf.word_wrap = True
        item_col  = WHITE if status == "DONE" else GRAY
        bullet_col = GREEN if status == "DONE" else YELLOW
        for item in items:
            add_para(tf, "    " + item, 12, color=item_col, before=5)

    add_transition(s)
    return s

# ── SLIDE 6: SYSTEM ARCHITECTURE ─────────────────────────────────────────────
def slide_architecture(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "System Architecture",
           "Firmware  ->  PicoRV32  ->  PCPI  ->  4x4 Systolic Accelerator  ->  Memory")
    footer(s)

    # Left: architecture layer summary
    rect(s, Inches(0.4), Inches(1.35), Inches(5.1), Inches(5.7),
         fill=CARD, line=CYAN, lw=Pt(1.5))
    box(s, "Architecture Layers",
        Inches(0.55), Inches(1.45), Inches(4.85), Inches(0.45),
        size=16, bold=True, color=CYAN)
    hrule(s, Inches(0.55), Inches(1.9), Inches(4.85), color=CYAN, thick=Pt(0.8))

    tb = s.shapes.add_textbox(Inches(0.55), Inches(2.0),
                               Inches(4.85), Inches(4.7))
    tf = tb.text_frame; tf.word_wrap = True
    layers = [
        ("Firmware (RISC-V C)",
         "Loads matrices at mapped addresses, issues custom opcode 0x5420818b", GREEN),
        ("PicoRV32 CPU",
         "RV32I core, ENABLE_PCPI=1, decodes custom opcode, stalls until pcpi_ready", CYAN),
        ("PCPI Wrapper",
         "Detects instruction, drives systolic engine, asserts pcpi_ready on finish", ORANGE),
        ("4x4 Systolic Array",
         "16 Q5.10 MAC units, output-stationary dataflow, 673 deterministic cycles", PINK),
        ("Flat Memory Map",
         "A: 0x100-0x13C   B: 0x140-0x17C   C (result): 0x200-0x23C", YELLOW),
    ]
    for title, desc, col in layers:
        add_para(tf, "  " + title, 14, bold=True, color=col, before=8)
        add_para(tf, "    " + desc, 11, color=WHITE, before=1)

    # Right: architecture diagram — white frame so the PNG blends cleanly
    rect(s, Inches(5.7), Inches(1.35), Inches(7.25), Inches(5.7),
         fill=WHITE, line=RGBColor(0xCC, 0xCC, 0xCC), lw=Pt(1))
    img(s, "tinyml_accelerator_detailed_architecture.png",
        Inches(5.75), Inches(1.4), Inches(7.15), Inches(5.6),
        subdir="diagram")

    add_transition(s)
    return s

# ── SLIDE 7: PCPI INTERFACE ───────────────────────────────────────────────────
def slide_pcpi(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "PCPI: The Custom Instruction Bridge",
           "Pico Co-Processor Interface -- zero-cost hardware dispatch")
    footer(s)

    # Left: signals table (compressed to top half)
    rect(s, Inches(0.4), Inches(1.35), Inches(5.8), Inches(3.05),
         fill=CARD, line=YELLOW, lw=Pt(1))
    box(s, "Interface Signals",
        Inches(0.55), Inches(1.45), Inches(5.5), Inches(0.4),
        size=14, bold=True, color=YELLOW)
    hrule(s, Inches(0.55), Inches(1.85), Inches(5.5), color=YELLOW, thick=Pt(0.8))

    signals = [
        ("pcpi_valid",  "CPU -> CoPro", "Custom instruction detected",        CYAN),
        ("pcpi_insn",   "CPU -> CoPro", "Full 32-bit instruction word",        CYAN),
        ("pcpi_rs1/2",  "CPU -> CoPro", "Source registers (matrix A/B addr)",  CYAN),
        ("pcpi_ready",  "CoPro -> CPU", "Computation done -- CPU resumes",     GREEN),
        ("pcpi_wr/rd",  "CoPro -> CPU", "Write-back flag + result (unused)",   ORANGE),
        ("pcpi_wait",   "CoPro -> CPU", "Hold CPU stall (optional)",           GRAY),
    ]
    tb = s.shapes.add_textbox(Inches(0.55), Inches(1.95), Inches(5.5), Inches(2.3))
    tf = tb.text_frame; tf.word_wrap = True
    for sig, direction, desc, col in signals:
        add_para(tf, f"  {sig}", 12, bold=True, color=col, before=3)
        add_para(tf, f"     {direction}  |  {desc}", 10, color=GRAY, before=0)

    # Right: handshake timeline (compressed)
    rect(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(3.05),
         fill=CARD2, line=CYAN, lw=Pt(1))
    box(s, "Execution Handshake",
        Inches(6.65), Inches(1.45), Inches(6.2), Inches(0.4),
        size=14, bold=True, color=CYAN)
    hrule(s, Inches(6.65), Inches(1.85), Inches(6.2), color=CYAN, thick=Pt(0.8))

    steps = [
        ("1", "CPU fetches  0x5420818b  custom opcode",            CYAN),
        ("2", "pcpi_valid=1, rs1/rs2 broadcast to coprocessor",    CYAN),
        ("3", "CPU stalls -- awaiting pcpi_ready=1",               YELLOW),
        ("4", "Coprocessor reads matrix A & B from memory",        ORANGE),
        ("5", "4x4 systolic array computes 16 dot products",       ORANGE),
        ("6", "Results written to matrix C region (0x200-0x23C)",  ORANGE),
        ("7", "pcpi_ready=1  -- CPU resumes next instruction",     GREEN),
    ]
    for i, (num, text, col) in enumerate(steps):
        y = Inches(1.95) + i * Inches(0.29)
        rect(s, Inches(6.65), y, Inches(0.3), Inches(0.28),
             fill=col, line=None)
        box(s, num, Inches(6.65), y, Inches(0.3), Inches(0.28),
            size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
        box(s, text, Inches(7.0), y + Inches(0.02), Inches(5.7), Inches(0.3),
            size=11, color=WHITE)

    # Opcode callout strip
    rect(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.38),
         fill=RGBColor(0x08, 0x0F, 0x1E), line=ORANGE, lw=Pt(1))
    box(s, "Custom opcode:  0x5420818b   (custom-0 space, funct7=0101010, rs2=00000, funct3=000, opcode=0001011)",
        Inches(0.55), Inches(4.54), Inches(12.1), Inches(0.3),
        size=11, bold=True, color=ORANGE)

    # PCPI waveform image
    img(s, "pcpi_integration_wave.png",
        Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.05))

    add_transition(s)
    return s

# ── SLIDE 8: SYSTOLIC ARRAY ───────────────────────────────────────────────────
def slide_systolic(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "The 4x4 Systolic Array",
           "16 parallel Processing Elements -- output-stationary dataflow")
    footer(s)

    # Left: systolic array diagram image
    img(s, "fig_systolic_array.png",
        Inches(0.4), Inches(1.35), Inches(6.5), Inches(5.7))

    # Right panel: dataflow explanation
    rx  = Inches(7.15)
    rw  = W - rx - Inches(0.2)
    rect(s, rx, Inches(1.35), rw, Inches(5.7),
         fill=CARD, line=CYAN, lw=Pt(1))
    box(s, "How It Works",
        rx + Inches(0.15), Inches(1.45), rw - Inches(0.2), Inches(0.45),
        size=16, bold=True, color=CYAN)
    hrule(s, rx + Inches(0.15), Inches(1.9), rw - Inches(0.2),
          color=CYAN, thick=Pt(0.8))

    tb = s.shapes.add_textbox(rx + Inches(0.15), Inches(2.0),
                               rw - Inches(0.2), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    facts = [
        ("Output-stationary",
         "Each PE accumulates its partial sum for C[i,j] in place"),
        ("Row broadcast",
         "Row i of matrix A flows right across row i of PEs"),
        ("Column broadcast",
         "Column j of matrix B flows down column j of PEs"),
        ("Parallel MACs",
         "All 16 PEs fire simultaneously every clock cycle"),
        ("Deterministic",
         "673 cycles regardless of operand values"),
        ("16-bit Q5.10",
         "Multiply -> shift-right 10 -> accumulate -> truncate"),
    ]
    for title, desc in facts:
        add_para(tf, "  " + title, 14, bold=True, color=ORANGE, before=8)
        add_para(tf, "    " + desc, 12, color=WHITE, before=1)

    add_transition(s)
    return s

# ── SLIDE 9: Q5.10 FIXED-POINT ────────────────────────────────────────────────
def slide_fixedpoint(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Q5.10 Fixed-Point Arithmetic",
           "16-bit two's complement -- integer simplicity, fractional precision")
    footer(s)

    # Bit layout diagram
    bit_w = Inches(0.78)
    bit_t = Inches(1.45)
    bit_h = Inches(0.75)

    labels   = ["S",
                "I4", "I3", "I2", "I1", "I0",
                "F9", "F8", "F7", "F6", "F5", "F4", "F3", "F2", "F1", "F0"]
    colors_  = [RED] + [ORANGE] * 5 + [CYAN] * 10
    bit_nums = [15, 14, 13, 12, 11, 10, 9, 8, 7, 6, 5, 4, 3, 2, 1, 0]

    for i in range(16):
        x = Inches(0.4) + i * bit_w
        rect(s, x, bit_t, bit_w, bit_h,
             fill=colors_[i], line=BG, lw=Pt(2))
        box(s, labels[i], x, bit_t + Inches(0.13), bit_w, bit_h - Inches(0.26),
            size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
        box(s, str(bit_nums[i]), x, bit_t + bit_h + Inches(0.02),
            bit_w, Inches(0.28), size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Brace labels
    box(s, "Sign (1 bit)",
        Inches(0.4), bit_t + bit_h + Inches(0.3),
        bit_w, Inches(0.38), size=11, color=RED, align=PP_ALIGN.CENTER)
    box(s, "Integer part  (5 bits)",
        Inches(1.2), bit_t + bit_h + Inches(0.3),
        bit_w * 5, Inches(0.38), size=11, color=ORANGE, align=PP_ALIGN.CENTER)
    box(s, "Fractional part  (10 bits)",
        Inches(1.2) + bit_w * 5, bit_t + bit_h + Inches(0.3),
        bit_w * 10, Inches(0.38), size=11, color=CYAN, align=PP_ALIGN.CENTER)

    # Properties grid
    props = [
        ("Range",      "-32.0  to  +31.999",            WHITE),
        ("Resolution", "1 / 1024  ~  0.001",            WHITE),
        ("Format",     "16-bit signed integer",         WHITE),
        ("Scaling",    "Multiply -> shift right 10 bits", WHITE),
    ]
    for i, (k, v, col) in enumerate(props):
        x = Inches(0.4) + i * Inches(3.22)
        rect(s, x, Inches(3.3), Inches(3.1), Inches(1.0),
             fill=CARD, line=CYAN, lw=Pt(1))
        box(s, k, x + Inches(0.1), Inches(3.38), Inches(2.9), Inches(0.38),
            size=13, bold=True, color=CYAN)
        box(s, v, x + Inches(0.1), Inches(3.72), Inches(2.9), Inches(0.48),
            size=14, color=WHITE)

    # PE arithmetic walkthrough
    rect(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.3),
         fill=CARD2, line=ORANGE, lw=Pt(1))
    box(s, "Processing Element MAC Operation  (per clock cycle):",
        Inches(0.55), Inches(4.58), Inches(12.1), Inches(0.42),
        size=15, bold=True, color=ORANGE)

    steps = [
        ("(1) Multiply",
         "a_in x b_in  ->  32-bit signed product"),
        ("(2) Scale",
         "product >>> 10  (arithmetic right-shift corrects Q5.10)"),
        ("(3) Accumulate",
         "acc +=  scaled_product  (32-bit running sum)"),
        ("(4) Pass through",
         "a_out <- a_in  ;  b_out <- b_in  (feed neighbours)"),
    ]
    for i, (step, desc) in enumerate(steps):
        x = Inches(0.55) + i * Inches(3.12)
        box(s, step, x, Inches(5.05), Inches(2.9), Inches(0.38),
            size=13, bold=True, color=YELLOW)
        box(s, desc, x, Inches(5.42), Inches(3.0), Inches(1.2),
            size=11, color=WHITE)

    add_transition(s)
    return s

# ── SLIDE 10: VERIFICATION STRATEGY ──────────────────────────────────────────
def slide_verification(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Verification Strategy",
           "Four-layer simulation campaign -- unit tests to full system benchmarking")
    footer(s)

    layers = [
        (4, "Performance Benchmarking",
         "Cycle-count measurement against two software baselines "
         "(no-MUL and MUL-equipped).",
         YELLOW, "673 cycles"),
        (3, "Protocol Handoff Validation",
         "Isolated test of PCPI signal handshake -- "
         "pcpi_valid, pcpi_ready, pcpi_insn correctness.",
         ORANGE, "PASS"),
        (2, "CPU Integration Regression",
         "Full system simulation: PicoRV32 + PCPI wrapper + accelerator. "
         "8 input profiles.",
         CYAN, "8 / 8 PASS"),
        (1, "Standalone Unit Tests",
         "Accelerator module in isolation. Zero, identity, dense, "
         "boundary, negative-element cases.",
         GREEN, "19 / 19 PASS"),
    ]

    for i, (num, title, desc, col, result) in enumerate(layers):
        y = Inches(1.35) + i * Inches(1.38)
        w_bar = Inches(12.5) - i * Inches(0.6)
        x_bar = Inches(0.4) + i * Inches(0.3)

        rect(s, x_bar, y, w_bar, Inches(1.28), fill=CARD, line=col, lw=Pt(1.5))

        rect(s, x_bar + Inches(0.1), y + Inches(0.19),
             Inches(0.55), Inches(0.55), fill=col)
        box(s, str(num), x_bar + Inches(0.1), y + Inches(0.19),
            Inches(0.55), Inches(0.55), size=20, bold=True,
            color=BG, align=PP_ALIGN.CENTER)

        box(s, f"Layer {num}  |  {title}",
            x_bar + Inches(0.8), y + Inches(0.1), Inches(9.0), Inches(0.48),
            size=17, bold=True, color=col)
        box(s, desc,
            x_bar + Inches(0.8), y + Inches(0.58), Inches(8.8), Inches(0.62),
            size=13, color=WHITE)

        # Result badge
        rect(s, x_bar + w_bar - Inches(2.0), y + Inches(0.35),
             Inches(1.85), Inches(0.5), fill=col)
        box(s, result,
            x_bar + w_bar - Inches(2.0), y + Inches(0.39),
            Inches(1.85), Inches(0.42), size=14, bold=True,
            color=BG, align=PP_ALIGN.CENTER)

    add_transition(s)
    return s

# ── SLIDE 11: TEST RESULTS SCORECARD ─────────────────────────────────────────
def slide_results_tests(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Verification Results",
           "100% pass rate across all test flows -- simulation-complete milestone")
    footer(s)

    # Big pass badge
    rect(s, Inches(0.4), Inches(1.35), Inches(3.8), Inches(5.7),
         fill=RGBColor(0x00, 0x20, 0x10), line=GREEN, lw=Pt(2))
    box(s, "100%", Inches(0.4), Inches(1.9), Inches(3.8), Inches(1.2),
        size=64, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    box(s, "PASS RATE",
        Inches(0.4), Inches(3.1), Inches(3.8), Inches(0.5),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hrule(s, Inches(0.7), Inches(3.65), Inches(3.2), color=GREEN, thick=Pt(1))

    tb = s.shapes.add_textbox(Inches(0.55), Inches(3.75), Inches(3.5), Inches(2.9))
    tf = tb.text_frame; tf.word_wrap = True
    for val, lbl in [("19 / 19", "standalone unit tests"),
                     ("8 / 8",   "integration cases"),
                     ("5 / 5",   "professor demo cases"),
                     ("PASS",    "handoff protocol test"),
                     ("PASS",    "one-command gate check")]:
        add_para(tf, f"{val}  {lbl}", 14, before=7,
                 color=GREEN if "19" in val or "8 /" in val else WHITE)

    # Test category breakdown
    categories = [
        ("Zero / Null Matrix",           2,  2,  GREEN),
        ("Identity Matrix",              3,  3,  GREEN),
        ("Dense (varied magnitudes)",    8,  8,  GREEN),
        ("Q5.10 Sign Boundary",          4,  4,  GREEN),
        ("Negative Elements",            2,  2,  GREEN),
        ("Integration: varied inputs",   8,  8,  CYAN),
        ("Professor Demo Cases",         5,  5,  ORANGE),
    ]
    rect(s, Inches(4.45), Inches(1.35), Inches(8.45), Inches(5.7),
         fill=CARD, line=CYAN, lw=Pt(1))
    box(s, "Test Category Breakdown",
        Inches(4.6), Inches(1.45), Inches(8.1), Inches(0.45),
        size=16, bold=True, color=CYAN)
    hrule(s, Inches(4.6), Inches(1.9), Inches(8.1), color=CYAN, thick=Pt(0.8))

    bar_l     = Inches(7.2)
    bar_max_w = Inches(4.8)
    for i, (cat, total, passed, col) in enumerate(categories):
        y = Inches(2.05) + i * Inches(0.68)
        box(s, cat, Inches(4.6), y + Inches(0.1), Inches(2.5), Inches(0.42),
            size=12, color=WHITE)
        rect(s, bar_l, y + Inches(0.18), bar_max_w, Inches(0.28),
             fill=CARD2, line=None)
        rect(s, bar_l, y + Inches(0.18),
             bar_max_w * passed / total, Inches(0.28), fill=col, line=None)
        box(s, f"{passed}/{total}", bar_l + bar_max_w + Inches(0.1),
            y + Inches(0.1), Inches(0.7), Inches(0.42),
            size=12, bold=True, color=col)

    add_transition(s)
    return s

# ── SLIDE 12: PERFORMANCE -- CYCLE COUNTS ─────────────────────────────────────
def slide_cycles(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Performance: Cycle Counts",
           "Deterministic execution -- 673 cycles regardless of input values")
    footer(s)

    # Left: number callouts
    rect(s, Inches(0.4), Inches(1.35), Inches(4.5), Inches(5.7),
         fill=CARD2, line=CYAN, lw=Pt(1.5))

    for i, (val, lbl, sub, col) in enumerate([
        ("673",    "EdgeMATX Accelerator",     "rv32i + PCPI",         CYAN),
        ("7,975",  "Software + MUL",           "rv32im baseline",      YELLOW),
        ("26,130", "Software (no MUL)",        "rv32i baseline",       RED),
        ("36,246", "Software (no MUL, dense)", "dense input profile",  RED),
    ]):
        y = Inches(1.45) + i * Inches(1.38)
        box(s, val,
            Inches(0.6), y, Inches(3.5), Inches(0.72),
            size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
        box(s, lbl,
            Inches(0.6), y + Inches(0.68), Inches(3.5), Inches(0.38),
            size=13, color=WHITE, align=PP_ALIGN.CENTER)
        box(s, sub,
            Inches(0.6), y + Inches(1.0), Inches(3.5), Inches(0.3),
            size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Right: bar chart image — white frame to match matplotlib background
    rect(s, Inches(5.1), Inches(1.35), Inches(7.9), Inches(5.7),
         fill=WHITE, line=RGBColor(0xCC, 0xCC, 0xCC), lw=Pt(1))
    img(s, "cycle_count_bar.png",
        Inches(5.15), Inches(1.4), Inches(7.8), Inches(5.6))

    add_transition(s)
    return s

# ── SLIDE 13: PERFORMANCE -- SPEEDUP ──────────────────────────────────────────
def slide_speedup(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Performance: Speedup",
           "38-54x over rv32i baseline -- 11.85x over MUL-equipped core")
    footer(s)

    # Top three callout cards
    for i, (val, title, desc, col) in enumerate([
        ("38.83x", "Sparse  vs  rv32i",
         "Identity-sequence benchmark\n26,130 -> 673 cycles",    GREEN),
        ("53.86x", "Dense   vs  rv32i",
         "Live dense profile\n36,246 -> 673 cycles",             CYAN),
        ("11.85x", "Any     vs  rv32im",
         "Even against MUL-equipped core\n7,975 -> 673 cycles",  ORANGE),
    ]):
        x = Inches(0.4) + i * Inches(4.32)
        rect(s, x, Inches(1.35), Inches(4.1), Inches(2.9),
             fill=CARD, line=col, lw=Pt(2))
        box(s, val, x, Inches(1.5), Inches(4.1), Inches(1.0),
            size=46, bold=True, color=col, align=PP_ALIGN.CENTER)
        box(s, title, x, Inches(2.5), Inches(4.1), Inches(0.45),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        box(s, desc, x, Inches(2.98), Inches(4.1), Inches(0.8),
            size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Bottom: speedup comparison chart — white frame
    rect(s, Inches(0.4), Inches(4.45), Inches(12.5), Inches(2.6),
         fill=WHITE, line=RGBColor(0xCC, 0xCC, 0xCC), lw=Pt(1))
    img(s, "speedup_comparison.png",
        Inches(0.45), Inches(4.5), Inches(12.4), Inches(2.5))

    add_transition(s)
    return s

# ── SLIDE 14: INTERACTIVE VISUALIZERS ─────────────────────────────────────────
def slide_visualizers(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Interactive Visualizers",
           "Explore the design live in your browser -- click links below")
    footer(s)

    vis_data = [
        ("Architecture Overview",
         "https://riscv-visualizer-2.vercel.app",
         "riscv-visualizer-2.vercel.app",
         ["High-level system architecture and dataflow",
          "Firmware layer  |  PicoRV32  |  PCPI  |  Systolic Array",
          "Step-through execution with annotations"],
         CYAN,
         "tinyml_accelerator_detailed_architecture.png",
         "diagram"),
        ("Processor Signal Simulation",
         "https://tinyml-pcpi-visualizer.vercel.app",
         "tinyml-pcpi-visualizer.vercel.app",
         ["Instruction decode  |  PCPI signal timeline",
          "Accelerator invocation waveform",
          "Systolic array computation state"],
         ORANGE,
         "pcpi_integration_wave.png",
         "figures"),
    ]

    for i, (title, url_full, url_short, desc_lines, col,
            preview_img, prev_subdir) in enumerate(vis_data):
        x = Inches(0.4) + i * Inches(6.5)

        # Card background
        rect(s, x, Inches(1.35), Inches(6.2), Inches(5.75),
             fill=CARD, line=col, lw=Pt(2))

        # Title
        box(s, title,
            x + Inches(0.2), Inches(1.45), Inches(5.8), Inches(0.5),
            size=18, bold=True, color=col)
        hrule(s, x + Inches(0.2), Inches(1.95), Inches(5.8),
              color=col, thick=Pt(0.8))

        # URL badge with clickable hyperlink
        rect(s, x + Inches(0.2), Inches(2.05), Inches(5.8), Inches(0.45),
             fill=RGBColor(0x08, 0x0F, 0x1E), line=col, lw=Pt(1))
        url_box(s, "  " + url_short, url_full,
                x + Inches(0.25), Inches(2.13), Inches(5.65), Inches(0.35),
                size=13, col=col)

        # Description text -- fixed height so it does NOT overlap the image below
        tb = s.shapes.add_textbox(x + Inches(0.2), Inches(2.62),
                                   Inches(5.8), Inches(0.95))
        tf = tb.text_frame; tf.word_wrap = True
        for line in desc_lines:
            add_para(tf, "    " + line, 12, color=WHITE, before=3)

        # Preview image — white frame so white-bg PNGs blend properly
        rect(s, x + Inches(0.2), Inches(3.65), Inches(5.8), Inches(3.1),
             fill=WHITE, line=col, lw=Pt(0.5))
        img(s, preview_img,
            x + Inches(0.22), Inches(3.67), Inches(5.76), Inches(3.06),
            subdir=prev_subdir)

    add_transition(s)
    return s

# ── SLIDE 15: FUTURE WORK ─────────────────────────────────────────────────────
def slide_future(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "What's Next: Pynq-Z2 FPGA Deployment",
           "Simulation-proven design -> physical hardware closure")
    footer(s)

    stages = [
        ("DONE",    "Simulation",
         "673-cycle verified\nbehaviour on\nIcarus Verilog",    GREEN),
        ("NEXT",    "FPGA Synthesis",
         "Map to Pynq-Z2\nXilinx Zynq-7020\nPost-synthesis timing", YELLOW),
        ("FUTURE",  "Scale & Deploy",
         "On-board cycle counts\nAXI-BRAM memory\n8x8 array extension", GRAY),
    ]
    for i, (status, stage, desc, col) in enumerate(stages):
        x = Inches(0.5) + i * Inches(4.28)
        rect(s, x, Inches(1.35), Inches(4.0), Inches(3.8),
             fill=CARD, line=col, lw=Pt(2))
        box(s, status, x, Inches(1.45), Inches(4.0), Inches(0.45),
            size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        box(s, stage, x, Inches(1.88), Inches(4.0), Inches(0.65),
            size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        box(s, desc, x, Inches(2.58), Inches(4.0), Inches(1.4),
            size=14, color=GRAY, align=PP_ALIGN.CENTER)

        if i < 2:
            box(s, "->",
                Inches(4.55) + i * Inches(4.28), Inches(2.75),
                Inches(0.4), Inches(0.5),
                size=26, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Open items panel
    rect(s, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.85),
         fill=CARD2, line=ORANGE, lw=Pt(1))
    box(s, "Open Items Before Hardware Closure",
        Inches(0.55), Inches(5.38), Inches(12.1), Inches(0.42),
        size=14, bold=True, color=ORANGE)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(5.82),
                               Inches(12.1), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    for item in [
        "Cycle-scaling estimator anchor needs updating (869 -> 673)",
        "On-board rdcycle CSR measurements to cross-validate simulation results",
        "AXI-BRAM memory interface replacing software-driven memory reads",
    ]:
        add_para(tf, "    " + item, 13, color=WHITE, before=3)

    add_transition(s)
    return s

# ── SLIDE 16: CONCLUSION ──────────────────────────────────────────────────────
def slide_conclusion(prs):
    s = blank(prs)
    set_bg(s)
    header(s, "Conclusion",
           "Simulation-complete milestone -- all objectives achieved")
    footer(s)

    achievements = [
        ("Hardware Accelerator",
         "4x4 systolic array with 16 Q5.10 fixed-point MACs, implemented in Verilog",
         CYAN),
        ("CPU Integration",
         "Tightly coupled to PicoRV32 via PCPI -- single opcode dispatch, zero overhead",
         ORANGE),
        ("Deterministic Execution",
         "673-cycle latency, invariant to operand values -- ideal for real-time systems",
         GREEN),
        ("Verified Performance",
         "38.83x - 53.86x speedup over rv32i;  11.85x over MUL-equipped rv32im",
         YELLOW),
        ("100% Test Coverage",
         "19/19 standalone  |  8/8 integration  |  5/5 professor demo  |  handoff PASS",
         PINK),
        ("Interactive Tooling",
         "Two live web visualizers for architecture walkthrough and signal-level simulation",
         GRAY),
    ]

    for i, (title, desc, col) in enumerate(achievements):
        row = i // 2
        c   = i %  2
        x = Inches(0.4) + c * Inches(6.3)
        y = Inches(1.35) + row * Inches(1.82)
        rect(s, x, y, Inches(6.05), Inches(1.66),
             fill=CARD, line=col, lw=Pt(1.5))
        box(s, title,
            x + Inches(0.15), y + Inches(0.12),
            Inches(5.75), Inches(0.48), size=16, bold=True, color=col)
        box(s, desc,
            x + Inches(0.15), y + Inches(0.64),
            Inches(5.75), Inches(0.85), size=12, color=WHITE)

    add_transition(s)
    return s

# ── SLIDE 17: THANK YOU ───────────────────────────────────────────────────────
def slide_thankyou(prs):
    s = blank(prs)
    set_bg(s)
    rect(s, Inches(0), Inches(0), W, Inches(0.08), fill=CYAN)
    rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=ORANGE)

    box(s, "Thank You",
        Inches(0.5), Inches(0.8), W - Inches(1.0), Inches(1.6),
        size=72, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    box(s, "Questions & Discussion",
        Inches(0.5), Inches(2.4), W - Inches(1.0), Inches(0.7),
        size=26, color=WHITE, align=PP_ALIGN.CENTER)

    hrule(s, Inches(2.0), Inches(3.2), Inches(9.333), color=ORANGE, thick=Pt(1.5))

    # Team
    tb = s.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(6.0), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    add_para(tf, "Team", 13, italic=True, color=GRAY)
    for name in ["Nishchay Pallav  221EC233",
                 "Mohammad Omar Sulemani  221EC230",
                 "Md Atib Kaif  221EC129"]:
        add_para(tf, name, 16, bold=True, color=WHITE, before=3)

    add_para(tf, "", 10, color=GRAY, before=6)
    add_para(tf, "Guide:  Dr Rathamala Rao", 13, italic=True, color=GRAY)

    # Live visualizers with hyperlinks
    tb2 = s.shapes.add_textbox(Inches(7.8), Inches(3.4), Inches(5.2), Inches(2.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    add_para(tf2, "Live Visualizers", 13, italic=True, color=GRAY)
    for display, full_url in [
        ("riscv-visualizer-2.vercel.app",       "https://riscv-visualizer-2.vercel.app"),
        ("tinyml-pcpi-visualizer.vercel.app",   "https://tinyml-pcpi-visualizer.vercel.app"),
    ]:
        p = tf2.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = display
        run.font.size = Pt(13)
        run.font.color.rgb = CYAN
        run.font.underline = True
        try:
            run.hyperlink.address = full_url
        except Exception:
            pass

    box(s,
        "Dept. of Electronics & Communication Engineering\n"
        "National Institute of Technology Karnataka, Surathkal  575025  |  2026",
        Inches(0.5), Inches(6.5), W - Inches(1.0), Inches(0.7),
        size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_transition(s)
    return s

# =============================================================================
# MAIN
# =============================================================================

def build():
    prs = make_prs()
    print("Building EdgeMATX presentation (17 slides)...")

    slide_title(prs)        ; print("  [OK] Slide  1 - Title")
    slide_introduction(prs) ; print("  [OK] Slide  2 - Introduction")
    slide_problem(prs)      ; print("  [OK] Slide  3 - The Problem")
    slide_solution(prs)     ; print("  [OK] Slide  4 - The Solution")
    slide_progress(prs)     ; print("  [OK] Slide  5 - Project Progress")
    slide_architecture(prs) ; print("  [OK] Slide  6 - System Architecture")
    slide_pcpi(prs)         ; print("  [OK] Slide  7 - PCPI Interface")
    slide_systolic(prs)     ; print("  [OK] Slide  8 - Systolic Array")
    slide_fixedpoint(prs)   ; print("  [OK] Slide  9 - Q5.10 Fixed-Point")
    slide_verification(prs) ; print("  [OK] Slide 10 - Verification Strategy")
    slide_results_tests(prs); print("  [OK] Slide 11 - Test Results")
    slide_cycles(prs)       ; print("  [OK] Slide 12 - Cycle Counts")
    slide_speedup(prs)      ; print("  [OK] Slide 13 - Speedup")
    slide_visualizers(prs)  ; print("  [OK] Slide 14 - Visualizers")
    slide_future(prs)       ; print("  [OK] Slide 15 - Future Work")
    slide_conclusion(prs)   ; print("  [OK] Slide 16 - Conclusion")
    slide_thankyou(prs)     ; print("  [OK] Slide 17 - Thank You")

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "EdgeMATX_presentation.pptx")
    prs.save(out)
    print("\n[OK] Done.  Saved -> " + out)

if __name__ == "__main__":
    build()
